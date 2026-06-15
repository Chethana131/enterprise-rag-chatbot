# ingest.py — Robust Semantic Indexer with Excel Support

import os
import json
import hashlib
import chromadb
import pandas as pd
import numpy as np
from sentence_transformers import SentenceTransformer
from PyPDF2 import PdfReader
from pptx import Presentation

DB_PATH = "./vector_db"
COLLECTION_NAME = "project_docs"
EMBED_MODEL = "all-MiniLM-L6-v2"

TEXT_EXT = {".txt", ".md", ".csv", ".xlsx", ".xls", ".pdf", ".ppt", ".pptx"}
BINARY_SKIP = {".h5", ".pt", ".pth", ".ckpt", ".onnx", ".npz", ".pkl"}


# =========================
# HELPERS
# =========================

def file_hash(path):
    h = hashlib.sha256()
    with open(path, "rb") as f:
        h.update(f.read())
    return h.hexdigest()


def safe_read_text(path):
    try:
        return open(path, "r", encoding="utf-8", errors="ignore").read()
    except:
        return ""


def make_columns_unique(cols):
    seen = {}
    new_cols = []

    for c in cols:
        c = str(c)
        if c not in seen:
            seen[c] = 0
            new_cols.append(c)
        else:
            seen[c] += 1
            new_cols.append(f"{c}_{seen[c]}")

    return new_cols


# =========================
# EXCEL TABLE DETECTION
# =========================

def detect_tables(df):

    tables = []
    start = None

    for i, row in df.iterrows():

        if row.notna().sum() > 1:
            if start is None:
                start = i
        else:
            if start is not None and i - start > 1:
                tables.append(df.iloc[start:i])
                start = None

    if start is not None:
        tables.append(df.iloc[start:])

    return tables


def detect_header(table):

    scores = []

    for i, row in table.iterrows():

        values = row.dropna()

        if len(values) == 0:
            scores.append(0)
            continue

        text_ratio = sum(isinstance(v, str) for v in values) / len(values)
        unique_ratio = len(set(values)) / len(values)

        scores.append(text_ratio + unique_ratio)

    return int(np.argmax(scores))


def clean_table(table):

    header_idx = detect_header(table)

    header = table.iloc[header_idx].astype(str).str.strip()

    data = table.iloc[header_idx + 1:].copy()

    data.columns = make_columns_unique(header)

    data = data.dropna(how="all")

    data = data.replace({np.nan: None})

    return data


def excel_tables_to_text(path):

    xls = pd.ExcelFile(path)

    texts = []
    structured = []

    for sheet in xls.sheet_names:

        raw = pd.read_excel(xls, sheet_name=sheet, header=None)

        raw = raw.replace({np.nan: None})

        tables = detect_tables(raw)

        if not tables:
            tables = [raw]

        for ti, table in enumerate(tables):

            try:

                df = clean_table(table)

                if df.empty:
                    continue

                preview = df.head(20)

                lines = [
                    f"Excel sheet={sheet} table={ti} columns={list(df.columns)}"
                ]

                for _, r in preview.iterrows():

                    row = [f"{c}={r.get(c)}" for c in df.columns]

                    lines.append(", ".join(row))

                texts.append("\n".join(lines))

                structured.append({
                    "sheet": sheet,
                    "table_id": ti,
                    "columns": list(df.columns),
                    "rows": preview.astype(str).to_dict(orient="records")
                })

            except Exception as e:

                print("Table parsing error:", e)

    return "\n\n".join(texts), structured


# =========================
# TEXT LOADING
# =========================

def load_text(path):

    ext = os.path.splitext(path)[1].lower()

    try:

        if ext == ".pdf":

            text = "\n".join(
                p.extract_text() or "" for p in PdfReader(path).pages
            )

            return text, None


        if ext == ".csv":

            df = pd.read_csv(path)

            return df.to_string(), None


        if ext in {".xlsx", ".xls"}:

            return excel_tables_to_text(path)


        if ext in {".ppt", ".pptx"}:

            prs = Presentation(path)

            text = "\n".join(
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )

            return text, None


        return safe_read_text(path), None


    except Exception as e:

        print("File read error:", e)

        return "", None


# =========================
# CHUNKING
# =========================

def chunk_text(text, max_len=900):

    chunks = []

    buf = ""

    for part in text.split("\n\n"):

        if len(buf) + len(part) < max_len:

            buf += part + "\n\n"

        else:

            chunks.append(buf.strip())

            buf = part

    if buf.strip():
        chunks.append(buf.strip())

    return chunks


# =========================
# INGEST
# =========================

def ingest(path):

    client = chromadb.PersistentClient(path=DB_PATH)

    col = client.get_or_create_collection(COLLECTION_NAME)

    embedder = SentenceTransformer(EMBED_MODEL)

    project = os.path.basename(os.path.abspath(path))

    items = []

    is_file = os.path.isfile(path)

    walker = os.walk(path) if os.path.isdir(path) else [
        (os.path.dirname(path), [], [os.path.basename(path)])
    ]

    for root, _, files in walker:

        for f in files:

            ext = os.path.splitext(f)[1].lower()

            if ext not in TEXT_EXT or ext in BINARY_SKIP:
                continue

            full = os.path.join(root, f)

            rel = f if is_file else os.path.relpath(full, path)

            print("Processing:", rel)

            raw, structured = load_text(full)

            if not raw.strip():
                print("No text extracted.")
                continue

            h = file_hash(full)

            # ---------- FILE METADATA (NEW FEATURE) ----------
            file_meta = {
                "file_name": os.path.basename(rel),
                "file_name_noext": os.path.splitext(os.path.basename(rel))[0],
                "file_path": rel,
                "file_type": ext
            }

            meta_text = json.dumps(file_meta)

            emb = embedder.encode([meta_text])[0]

            items.append((
                f"file_meta::{project}/{rel}",
                meta_text,
                emb,
                {
                    "type": "file_meta",
                    "file": rel,
                    "filename": os.path.basename(rel),
                    "filename_noext": os.path.splitext(os.path.basename(rel))[0],
                    "hash": h
                }
            ))
            # --------------------------------------------------

            chunks = chunk_text(raw)

            print("Chunks created:", len(chunks))

            for i, c in enumerate(chunks):

                if not c.strip():
                    continue

                emb = embedder.encode([c])[0]

                items.append((
                    f"text::{project}/{rel}::{i}",
                    c,
                    emb,
                    {
                        "type": "text",
                        "file": rel,
                        "filename": os.path.basename(rel),
                        "filename_noext": os.path.splitext(os.path.basename(rel))[0],
                        "hash": h
                    }
                ))

            if structured:

                meta = json.dumps(structured, default=str)

                emb = embedder.encode([meta])[0]

                items.append((
                    f"excel_struct::{project}/{rel}",
                    meta,
                    emb,
                    {
                        "type": "excel_struct",
                        "file": rel,
                        "filename": os.path.basename(rel),
                        "filename_noext": os.path.splitext(os.path.basename(rel))[0],
                        "hash": h
                    }
                ))

    if items:

        col.add(
            ids=[i[0] for i in items],
            documents=[i[1] for i in items],
            embeddings=[i[2].tolist() for i in items],
            metadatas=[i[3] for i in items]
        )

    print(f"\nIndexed {len(items)} embeddings")


if __name__ == "__main__":

    import argparse

    p = argparse.ArgumentParser()

    p.add_argument("path")

    args = p.parse_args()

    ingest(args.path)